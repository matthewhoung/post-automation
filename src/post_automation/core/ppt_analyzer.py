"""PowerPoint text extraction and analysis."""

from typing import List

from pptx import Presentation

from post_automation.models.ppt import SlideText
from post_automation.utils.logger import setup_logger

logger = setup_logger(__name__)


class PPTAnalyzer:
    """Analyzer for extracting text from PowerPoint presentations."""

    def extract_text_from_pptx(self, pptx_path: str) -> List[SlideText]:
        """
        Extract text from PowerPoint file.

        Args:
            pptx_path: Path to PPTX file.

        Returns:
            List of SlideText objects with extracted text.
        """
        logger.info(f"Extracting text from: {pptx_path}")

        try:
            prs = Presentation(pptx_path)
            results = []

            for i, slide in enumerate(prs.slides):
                slide_text_parts = []

                # Extract text from all shapes
                for shape in slide.shapes:
                    text = self._extract_text_from_shape(shape)
                    if text.strip():
                        slide_text_parts.append(text)

                # Combine all text parts
                combined_text = " ".join(slide_text_parts)

                results.append(
                    SlideText(
                        slide_number=i, text=combined_text, shape_count=len(slide.shapes)
                    )
                )

                logger.debug(
                    f"Slide {i}: Extracted {len(combined_text)} characters "
                    f"from {len(slide.shapes)} shapes"
                )

            logger.info(f"Successfully extracted text from {len(results)} slides")
            return results

        except Exception as e:
            logger.error(f"Failed to extract text from {pptx_path}: {e}")
            raise

    def _extract_text_from_shape(self, shape) -> str:
        """
        Extract text from a single shape.

        Args:
            shape: PowerPoint shape object.

        Returns:
            Extracted text.
        """
        text_parts = []

        # Handle text frames
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    text_parts.append(para_text)

        # Handle tables
        if hasattr(shape, "table"):
            table_text = self._extract_table_text(shape.table)
            if table_text:
                text_parts.append(table_text)

        # Handle grouped shapes (if present)
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                sub_text = self._extract_text_from_shape(sub_shape)
                if sub_text.strip():
                    text_parts.append(sub_text)

        return " ".join(text_parts)

    def _extract_table_text(self, table) -> str:
        """
        Extract text from a table.

        Args:
            table: PowerPoint table object.

        Returns:
            Combined table text.
        """
        text_parts = []

        try:
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        text_parts.append(cell_text)
        except Exception as e:
            logger.warning(f"Failed to extract table text: {e}")

        return " ".join(text_parts)

    def get_presentation_info(self, pptx_path: str) -> dict:
        """
        Get basic information about a presentation.

        Args:
            pptx_path: Path to PPTX file.

        Returns:
            Dictionary with presentation metadata.
        """
        try:
            prs = Presentation(pptx_path)

            return {
                "total_slides": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
            }

        except Exception as e:
            logger.error(f"Failed to get presentation info: {e}")
            raise
