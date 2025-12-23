"""PowerPoint modification utilities."""

from typing import List, Optional

from pptx import Presentation
from pptx.util import Pt, RGBColor

from post_automation.models.ppt import Replacement, StyleConfig
from post_automation.utils.logger import setup_logger

logger = setup_logger(__name__)


class PPTModifier:
    """Modifier for PowerPoint presentations."""

    def __init__(self, pptx_path: str):
        """
        Initialize PPT modifier.

        Args:
            pptx_path: Path to PowerPoint file.
        """
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)
        logger.info(f"Loaded presentation: {pptx_path}")

    def replace_content(self, replacements: List[Replacement]) -> None:
        """
        Replace text content in slides.

        Args:
            replacements: List of text replacements to apply.
        """
        logger.info(f"Applying {len(replacements)} content replacements")

        for replacement in replacements:
            try:
                self._replace_in_slide(
                    replacement.slide_num, replacement.old_text, replacement.new_text
                )
                logger.debug(
                    f"Replaced text in slide {replacement.slide_num}: "
                    f"'{replacement.old_text[:50]}...' -> '{replacement.new_text[:50]}...'"
                )
            except Exception as e:
                logger.error(
                    f"Failed to replace content in slide {replacement.slide_num}: {e}"
                )

    def _replace_in_slide(self, slide_num: int, old_text: str, new_text: str) -> None:
        """
        Replace text in a specific slide.

        Args:
            slide_num: Slide number (0-indexed).
            old_text: Text to replace.
            new_text: Replacement text.
        """
        if slide_num >= len(self.presentation.slides):
            logger.warning(f"Slide {slide_num} does not exist")
            return

        slide = self.presentation.slides[slide_num]

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                self._replace_text_in_frame(shape.text_frame, old_text, new_text)

            # Handle tables
            if hasattr(shape, "table"):
                self._replace_text_in_table(shape.table, old_text, new_text)

    def _replace_text_in_frame(self, text_frame, old_text: str, new_text: str) -> None:
        """
        Replace text in a text frame while preserving formatting.

        Args:
            text_frame: PowerPoint text frame.
            old_text: Text to replace.
            new_text: Replacement text.
        """
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    # Save original formatting
                    font_name = run.font.name
                    font_size = run.font.size
                    font_bold = run.font.bold
                    font_italic = run.font.italic
                    font_color = run.font.color.rgb if run.font.color.type == 1 else None

                    # Replace text
                    run.text = run.text.replace(old_text, new_text)

                    # Restore formatting
                    if font_name:
                        run.font.name = font_name
                    if font_size:
                        run.font.size = font_size
                    if font_bold is not None:
                        run.font.bold = font_bold
                    if font_italic is not None:
                        run.font.italic = font_italic
                    if font_color:
                        run.font.color.rgb = font_color

    def _replace_text_in_table(self, table, old_text: str, new_text: str) -> None:
        """
        Replace text in a table.

        Args:
            table: PowerPoint table.
            old_text: Text to replace.
            new_text: Replacement text.
        """
        try:
            for row in table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        self._replace_text_in_frame(cell.text_frame, old_text, new_text)
        except Exception as e:
            logger.warning(f"Failed to replace text in table: {e}")

    def modify_styles(self, style_config: StyleConfig) -> None:
        """
        Apply style changes to presentation.

        Args:
            style_config: Style configuration to apply.
        """
        logger.info("Applying style modifications")

        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    self._apply_styles_to_frame(shape.text_frame, style_config)

                # Handle tables
                if hasattr(shape, "table"):
                    self._apply_styles_to_table(shape.table, style_config)

    def _apply_styles_to_frame(self, text_frame, style_config: StyleConfig) -> None:
        """
        Apply styles to a text frame.

        Args:
            text_frame: PowerPoint text frame.
            style_config: Style configuration.
        """
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Apply font name
                if style_config.font_name:
                    run.font.name = style_config.font_name

                # Apply font size
                if style_config.font_size:
                    run.font.size = Pt(style_config.font_size)

                # Apply colors
                if style_config.color_scheme and style_config.color_scheme.text_color:
                    rgb = self._hex_to_rgb(style_config.color_scheme.text_color)
                    run.font.color.rgb = RGBColor(*rgb)

    def _apply_styles_to_table(self, table, style_config: StyleConfig) -> None:
        """
        Apply styles to a table.

        Args:
            table: PowerPoint table.
            style_config: Style configuration.
        """
        try:
            for row in table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        self._apply_styles_to_frame(cell.text_frame, style_config)
        except Exception as e:
            logger.warning(f"Failed to apply styles to table: {e}")

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """
        Convert hex color to RGB tuple.

        Args:
            hex_color: Hex color string (e.g., "#FF5733" or "FF5733").

        Returns:
            RGB tuple (r, g, b).
        """
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def save(self, output_path: str) -> None:
        """
        Save modified presentation.

        Args:
            output_path: Path where to save the presentation.
        """
        try:
            self.presentation.save(output_path)
            logger.info(f"Saved modified presentation to: {output_path}")
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise

    def get_slide_count(self) -> int:
        """
        Get number of slides in presentation.

        Returns:
            Number of slides.
        """
        return len(self.presentation.slides)
